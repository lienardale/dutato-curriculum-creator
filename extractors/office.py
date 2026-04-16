"""
Office document extractor — DOCX and PPTX to the unified intermediate format.

DOCX: Delegates to pdf_pipeline's extract.extract_docx() then converts.
PPTX: Uses python-pptx to extract slide content as sections.
"""

import sys
from pathlib import Path

import _compat  # noqa: F401


def _extract_docx_images(file_path: Path, images_dir: str) -> list[dict]:
    """Extract images from a DOCX file and write them to disk."""
    from docx import Document
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    out = Path(images_dir)
    out.mkdir(parents=True, exist_ok=True)

    doc = Document(str(file_path))
    registry: list[dict] = []
    img_idx = 0

    for rel in doc.part.rels.values():
        if "image" not in rel.reltype:
            continue
        try:
            blob = rel.target_part.blob
        except Exception:
            continue
        if len(blob) < 2048:
            continue

        content_type = rel.target_part.content_type or "image/png"
        ext = content_type.split("/")[-1].replace("jpeg", "jpg")
        img_id = f"docx_img{img_idx}"
        filename = f"{img_id}.{ext}"
        filepath = out / filename
        filepath.write_bytes(blob)

        registry.append({
            "id": img_id,
            "local_path": f"images/{filename}",
            "mime_type": content_type,
            "size_bytes": len(blob),
            "width": 0,
            "height": 0,
        })
        img_idx += 1

    return registry


def _extract_docx(file_path: Path, *, images_dir: str | None = None) -> dict:
    """Extract a DOCX file."""
    from extract import extract_docx as _raw_extract_docx

    raw = _raw_extract_docx(file_path)
    pages = raw["pages"]
    toc = raw["toc"]

    # Extract images
    image_registry: list[dict] = []
    if images_dir:
        image_registry = _extract_docx_images(file_path, images_dir)

    # Convert TOC entries to sections with content
    sections = []
    if toc:
        for i, entry in enumerate(toc):
            content_parts = []
            for page in pages:
                if entry["title"] in page["text"]:
                    content_parts.append(page["text"])
            sections.append({
                "title": entry["title"],
                "content": "\n\n".join(content_parts),
                "depth": entry["level"] - 1,
                "metadata": {},
            })
    else:
        full_text = "\n\n".join(p["text"] for p in pages)
        sections.append({
            "title": file_path.stem,
            "content": full_text,
            "depth": 0,
            "metadata": {},
        })

    # Distribute images evenly across sections (best-effort; DOCX lacks
    # reliable positional info for inline images)
    if image_registry and sections:
        per_section = max(1, len(image_registry) // len(sections))
        img_iter = iter(image_registry)
        for sec in sections:
            sec_images = []
            for _ in range(per_section):
                img = next(img_iter, None)
                if img:
                    sec_images.append({
                        "id": img["id"],
                        "local_path": img["local_path"],
                        "alt_text": "",
                        "context": "document",
                    })
            if sec_images:
                sec["images"] = sec_images
        # Attach remaining images to last section
        remaining = list(img_iter)
        if remaining:
            last = sections[-1]
            last.setdefault("images", []).extend(
                {"id": img["id"], "local_path": img["local_path"],
                 "alt_text": "", "context": "document"}
                for img in remaining
            )

    total_tokens = sum(len(s["content"].split()) for s in sections if s["content"])

    result: dict = {
        "source_type": "docx",
        "source_path": str(file_path.resolve()),
        "title": file_path.stem,
        "author": "",
        "sections": sections,
        "metadata": {
            "total_sections": len(sections),
            "total_tokens": total_tokens,
            "total_images": len(image_registry),
        },
    }
    if image_registry:
        result["images"] = image_registry
    return result


def _extract_pptx(file_path: Path, *, images_dir: str | None = None) -> dict:
    """Extract a PPTX file — each slide becomes a section."""
    from pptx import Presentation

    out = Path(images_dir) if images_dir else None
    if out:
        out.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(file_path))

    sections = []
    image_registry: list[dict] = []

    for i, slide in enumerate(prs.slides, 1):
        # Get slide title
        title = f"Slide {i}"
        if slide.shapes.title:
            title = slide.shapes.title.text.strip() or title

        # Collect all text and images from the slide
        texts = []
        slide_images: list[dict] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)

            # Extract images from shapes
            if out and hasattr(shape, "image"):
                try:
                    blob = shape.image.blob
                    if len(blob) < 2048:
                        continue
                    content_type = shape.image.content_type or "image/png"
                    ext = content_type.split("/")[-1].replace("jpeg", "jpg")
                    img_id = f"s{i}_img{len(slide_images)}"
                    filename = f"{img_id}.{ext}"
                    filepath = out / filename
                    filepath.write_bytes(blob)

                    img_entry = {
                        "id": img_id,
                        "local_path": f"images/{filename}",
                        "mime_type": content_type,
                        "size_bytes": len(blob),
                        "width": 0,
                        "height": 0,
                        "slide": i,
                    }
                    image_registry.append(img_entry)
                    slide_images.append({
                        "id": img_id,
                        "local_path": img_entry["local_path"],
                        "alt_text": "",
                        "context": f"slide:{i}",
                    })
                except Exception:
                    pass

        content = "\n\n".join(texts)

        # Collect notes if present
        notes = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes = notes_frame.text.strip()

        section: dict = {
            "title": title,
            "content": content,
            "depth": 0,
            "metadata": {
                "slide_number": i,
                "has_notes": bool(notes),
                "notes": notes,
            },
        }
        if slide_images:
            section["images"] = slide_images
        sections.append(section)

    total_tokens = sum(len(s["content"].split()) for s in sections if s["content"])

    result: dict = {
        "source_type": "pptx",
        "source_path": str(file_path.resolve()),
        "title": file_path.stem,
        "author": "",
        "sections": sections,
        "metadata": {
            "total_sections": len(sections),
            "total_tokens": total_tokens,
            "total_slides": len(prs.slides),
            "total_images": len(image_registry),
        },
    }
    if image_registry:
        result["images"] = image_registry
    return result


def extract_office(source: str, *, images_dir: str | None = None) -> dict:
    """Extract DOCX or PPTX file."""
    file_path = Path(source)
    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    suffix = file_path.suffix.lower()
    if suffix == ".docx":
        return _extract_docx(file_path, images_dir=images_dir)
    elif suffix == ".pptx":
        return _extract_pptx(file_path, images_dir=images_dir)
    else:
        raise ValueError(f"Unsupported office format: {suffix}")


if __name__ == "__main__":
    import json

    if len(sys.argv) < 2:
        print("Usage: python -m extractors.office <input.docx|pptx> [-o output_dir]")
        sys.exit(1)

    from extractors import extract_source
    source_path = sys.argv[1]
    output_dir = sys.argv[sys.argv.index("-o") + 1] if "-o" in sys.argv else None
    result = extract_source(source_path, output_dir)
    if not output_dir:
        print(json.dumps(result, indent=2, ensure_ascii=False))
    else:
        print(f"Extracted {result['metadata']['total_sections']} sections")
